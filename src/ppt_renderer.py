from __future__ import annotations

from io import BytesIO
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt

from .models import SheetTable

EMU_PER_INCH = 914400
EMU_PER_PX = 9525
DEFAULT_HEADER_LAYOUT_GAP_PX = 24
TITLE_P1 = "SINAIS DE COMPLICA\u00C7\u00C3O | RESULTADO POR ESTADO"
TITLE_P3 = "ORIENTACOES DO CUIDADO POS CIRURGICO | RESULTADO POR ESTADO"


def _inches_to_emu(value_in_inches: float) -> int:
    return int(Inches(value_in_inches))


def _emu_to_inches(value_emu: int) -> float:
    return float(value_emu) / EMU_PER_INCH


def _as_text(value: object, header: str) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        value = int(value)

    header_norm = header.upper()
    is_percent_col = (
        "PERCENTUAL" in header_norm
        or "PROPORCIONAL" in header_norm
        or "REPRESENTAT" in header_norm
    )
    if isinstance(value, (int, float)) and is_percent_col:
        num = float(value)
        if abs(num) <= 1:
            num *= 100
        return f"{num:.1f}%"
    if isinstance(value, int):
        return f"{value:,}".replace(",", ".")
    if isinstance(value, float):
        return f"{value:.1f}".replace(".", ",")
    return str(value)


def _is_general_sheet(name: str) -> bool:
    return name.upper() == "GERAL" or name.upper().endswith("_GERAL")


def _is_p3_sheet(name: str) -> bool:
    return name.upper().startswith("P3_")


def _title_for_sheet(name: str) -> str:
    return TITLE_P3 if _is_p3_sheet(name) else TITLE_P1


def _find_primary_general_sheet(sheet_names: list[str]) -> str:
    if "GERAL" in sheet_names:
        return "GERAL"
    preferred = [name for name in sheet_names if name.upper() == "P1_GERAL"]
    if preferred:
        return preferred[0]
    general_candidates = [name for name in sheet_names if name.upper().endswith("_GERAL")]
    if general_candidates:
        return general_candidates[0]
    raise ValueError('Nenhuma aba geral encontrada. Esperado "GERAL", "P1_GERAL" ou sufixo "_GERAL".')


def _set_slide_title(slide, text: str, slide_w: int, title_h: int, assets_dir: Path) -> None:
    side_margin = int(slide_w * 0.03)

    inicio_path = assets_dir.parent / "imgs" / "inicio.jpg"
    if not inicio_path.exists():
        inicio_path = assets_dir / "inicio.jpg"

    inicio_h = int(title_h * 0.70)
    inicio_w = int(inicio_h * (100 / 681))
    inicio_x = side_margin
    inicio_y = int((title_h - inicio_h) / 2)
    if inicio_path.exists():
        slide.shapes.add_picture(str(inicio_path), inicio_x, inicio_y, width=inicio_w, height=inicio_h)

    text_gap = int(20 * EMU_PER_PX)
    text_left = inicio_x + inicio_w + text_gap
    text_top = int(title_h * 0.08)
    text_w = int(slide_w - text_left - side_margin - int(slide_w * 0.10))
    text_h = int(title_h * 0.84)
    box = slide.shapes.add_textbox(text_left, text_top, text_w, text_h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.name = "Calibri Light"
    run.font.color.rgb = RGBColor(32, 56, 100)
    run.font.size = Pt(26)

    logo_path = assets_dir.parent / "imgs" / "logo.jpg"
    if not logo_path.exists():
        logo_path = assets_dir / "logo.jpg"
    logo_w = int(50 * EMU_PER_PX)
    logo_h = int(50 * EMU_PER_PX)
    logo_padding_right = int(25 * EMU_PER_PX)
    logo_padding_top = int(25 * EMU_PER_PX)
    logo_x = int(slide_w - logo_w - logo_padding_right)
    logo_y = logo_padding_top
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), logo_x, logo_y, width=logo_w, height=logo_h)


def _column_weights(headers: list[str], rows: list[list[object]]) -> list[float]:
    widths = []
    for c, h in enumerate(headers):
        max_len = len(str(h))
        for row in rows:
            if c < len(row):
                max_len = max(max_len, len(_as_text(row[c], h)))
        widths.append(max(1.0, float(max_len)))
    return widths


def _table_height_with_row_cap(available_h: int, row_count: int, slide_h: int) -> int:
    if row_count <= 0:
        return max(1, available_h)
    max_row_h = int(slide_h * 0.045)
    min_row_h = int(slide_h * 0.018)
    capped_h = min(available_h, row_count * max_row_h)
    floor_h = row_count * min_row_h
    return max(1, max(capped_h, floor_h))


def _calc_uniform_font_size_pt(
    headers: list[str],
    body: list[list[str]],
    col_widths: list[float],
    width_in: float,
    height_in: float,
) -> float:
    row_count = len(body) + 1
    available_pt_h = height_in * 72 * 0.92
    vertical_limit_pt = available_pt_h / max(2.0, row_count * 1.30)

    max_chars_per_col = []
    for c in range(len(headers)):
        max_len = len(str(headers[c]))
        for row in body:
            if c < len(row):
                max_len = max(max_len, len(str(row[c])))
        max_chars_per_col.append(max_len)

    horizontal_limits = []
    for c in range(len(headers)):
        col_width_pt = width_in * 72 * col_widths[c] * 0.96
        char_count = max(1, max_chars_per_col[c])
        horizontal_limits.append(col_width_pt / (char_count * 0.62))

    horizontal_limit_pt = min(horizontal_limits) if horizontal_limits else 8.0
    return max(3.0, min(7.2, vertical_limit_pt, horizontal_limit_pt))


def _render_table_image(
    table_data: SheetTable,
    width_in: float,
    height_in: float,
    dpi: int = 220,
) -> BytesIO:
    headers = table_data.headers or ["DADO"]
    rows = table_data.rows or []
    col_count = max(1, len(headers))
    body = [[_as_text(row[c] if c < len(row) else "", headers[c]) for c in range(col_count)] for row in rows]

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=dpi)
    fig.subplots_adjust(left=0.0, right=1.0, bottom=0.0, top=1.0)
    ax.axis("off")

    weights = _column_weights(headers, rows)
    weight_sum = sum(weights) if sum(weights) > 0 else 1.0
    col_widths = [w / weight_sum for w in weights]
    if col_count > 1:
        first_col = col_widths[0] * 1.08
        rest_sum = sum(col_widths[1:])
        if rest_sum > 0:
            scale = max(0.0, (1.0 - first_col) / rest_sum)
            col_widths = [first_col] + [w * scale for w in col_widths[1:]]

    table_font_size = _calc_uniform_font_size_pt(headers, body, col_widths, width_in, height_in)
    table = ax.table(
        cellText=body,
        colLabels=headers,
        colWidths=col_widths,
        cellLoc="center",
        bbox=[0, 0, 1, 1],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(table_font_size)

    uniform_row_h = 1.0 / max(1, (len(body) + 1))
    for r in range(len(body) + 1):
        for c in range(col_count):
            table[(r, c)].set_height(uniform_row_h)

    blue = "#2652B5"
    zebra = "#EDEDED"
    white = "#FFFFFF"
    black = "#000000"

    for (r, c), cell in table.get_celld().items():
        cell.set_linewidth(0.6)
        cell.set_edgecolor(white)

        txt = cell.get_text()
        txt.set_fontname("Calibri")
        txt.set_fontweight("bold")
        txt.set_linespacing(1.0)
        txt.set_va("center")
        txt.set_wrap(False)
        if r == 0 or c == 0:
            txt.set_fontsize(max(2.6, table_font_size * 0.86))
        else:
            txt.set_fontsize(table_font_size)

        if r == 0:
            cell.set_facecolor(blue)
            txt.set_color(white)
            txt.set_ha("center")
            continue

        data_idx = r - 1
        first_value = rows[data_idx][0] if data_idx < len(rows) and len(rows[data_idx]) > 0 else ""
        is_total = str(first_value).strip().upper() == "TOTAL"

        if is_total:
            cell.set_facecolor(blue)
            txt.set_color(white)
            txt.set_ha("left" if c == 0 else "center")
        elif c == 0:
            cell.set_facecolor(blue)
            txt.set_color(white)
            txt.set_ha("left")
        else:
            cell.set_facecolor(zebra if r % 2 == 0 else white)
            txt.set_color(black)
            txt.set_ha("center")

    buffer = BytesIO()
    fig.savefig(buffer, format="png", dpi=dpi, facecolor="white", bbox_inches="tight", pad_inches=0)
    plt.close(fig)
    buffer.seek(0)
    return buffer


def _add_table_block(
    slide,
    table_data: SheetTable,
    block_left: int,
    block_top: int,
    block_w: int,
    block_h: int,
    slide_h: int,
    table_title_h: int | None = None,
) -> None:
    title_h_for_layout = table_title_h if table_title_h is not None else int(block_h * 0.20)
    title_bottom_gap_for_layout = int(block_h * 0.03)

    row_count = len(table_data.rows) + 1
    table_top = block_top + title_h_for_layout + title_bottom_gap_for_layout
    available_table_h = max(1, block_h - title_h_for_layout - title_bottom_gap_for_layout)
    table_h = _table_height_with_row_cap(available_table_h, row_count, slide_h)
    table_top += int((available_table_h - table_h) / 2)

    if table_data.display_name.strip().upper() != "GERAL":
        title_gap = int(12 * EMU_PER_PX)
        title_h_visual = int(30 * EMU_PER_PX)
        title_top = max(0, table_top - title_gap - title_h_visual)
        title_box = slide.shapes.add_textbox(block_left, title_top, block_w, title_h_visual)
        tf = title_box.text_frame
        tf.clear()
        tf.margin_left = int(50 * EMU_PER_PX)
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = table_data.display_name
        run.font.bold = True
        run.font.color.rgb = RGBColor(31, 56, 100)
        run.font.size = Pt(20)

    width_in = _emu_to_inches(block_w)
    height_in = _emu_to_inches(table_h)
    image_buffer = _render_table_image(table_data, width_in=width_in, height_in=height_in)
    slide.shapes.add_picture(image_buffer, block_left, table_top, width=block_w, height=table_h)


def _render_pie_image(
    title: str,
    sim_value: int,
    nao_value: int,
    width_in: float,
    height_in: float,
    dpi: int = 220,
) -> BytesIO:
    labels = ["Sim", "Não"]
    values = [int(sim_value), int(nao_value)]
    colors = ["#2652B5", "#E64A36"]

    total = sum(values)
    if total <= 0:
        values = [0, 0]
        total = 0

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=dpi)
    fig.patch.set_facecolor("#FFFFFF")
    ax.set_facecolor("#FFFFFF")

    wedges, _texts = ax.pie(
        values,
        labels=None,
        colors=colors,
        startangle=90,
        counterclock=False,
        radius=0.86,
        wedgeprops={"linewidth": 0.0},
    )

    for wedge, value in zip(wedges, values):
        theta = (wedge.theta1 + wedge.theta2) / 2.0
        x = float(math.cos(math.radians(theta)))
        y = float(math.sin(math.radians(theta)))
        label_x = 1.15 * x
        label_y = 1.15 * y
        pct = 0.0 if total == 0 else (value / total * 100.0)
        text = f"{value} ({pct:.0f}%)"
        ha = "left" if x >= 0 else "right"
        ax.annotate(
            text,
            xy=(0.86 * x, 0.86 * y),
            xytext=(label_x, label_y),
            ha=ha,
            va="center",
            fontsize=12,
            fontname="Calibri",
            fontweight="bold",
            color="#1A1A1A",
            arrowprops={"arrowstyle": "-", "color": "#B3B3B3", "linewidth": 1.0},
        )

    ax.set_aspect("equal")
    title_band = plt.Rectangle(
        (0.0, 0.92),
        1.00,
        0.08,
        transform=fig.transFigure,
        facecolor="#2652B5",
        edgecolor="none",
        clip_on=False,
        zorder=3,
    )
    fig.add_artist(title_band)
    fig.text(
        0.02,
        0.955,
        title,
        ha="left",
        va="center",
        fontsize=15,
        fontname="Calibri",
        fontweight="bold",
        color="#FFFFFF",
        zorder=4,
    )
    ax.legend(
        wedges,
        labels,
        loc="lower center",
        bbox_to_anchor=(0.5, -0.10),
        ncol=2,
        frameon=False,
        handlelength=0.9,
        handletextpad=0.3,
        columnspacing=1.0,
        prop={"family": "Calibri", "weight": "bold", "size": 12},
    )

    buffer = BytesIO()
    fig.savefig(buffer, format="png", dpi=dpi, facecolor=fig.get_facecolor(), pad_inches=0)
    plt.close(fig)
    buffer.seek(0)
    return buffer


def _add_pie_block(
    slide,
    title: str,
    sim_value: int,
    nao_value: int,
    block_left: int,
    block_top: int,
    block_w: int,
    block_h: int,
) -> None:
    width_in = _emu_to_inches(block_w)
    height_in = _emu_to_inches(block_h)
    image_buffer = _render_pie_image(
        title=title,
        sim_value=sim_value,
        nao_value=nao_value,
        width_in=width_in,
        height_in=height_in,
    )
    slide.shapes.add_picture(image_buffer, block_left, block_top, width=block_w, height=block_h)


def _add_requested_pies(
    prs: Presentation,
    slide_w: int,
    slide_h: int,
    contagens_sim_nao: dict[str, dict[str, int]] | None,
    header_layout_gap_emu: int,
) -> None:
    if not contagens_sim_nao:
        return

    page_title_h = int(slide_h * 0.20)
    outer_margin = int(slide_w * 0.04)
    center_gap = int(slide_w * 0.03)
    panel_w = int((slide_w - (2 * outer_margin) - center_gap) / 2)
    available_h = slide_h - page_title_h - header_layout_gap_emu
    panel_h = int(available_h * 0.78)
    panel_top = page_title_h + header_layout_gap_emu + int((available_h - panel_h) / 2)
    left_x = outer_margin
    right_x = outer_margin + panel_w + center_gap

    pie_specs = [
        (3, "P1", "P1. Você percebeu inchaço ou caroço no corte da cirurgia?"),
        (15, "P3", "P3"),
    ]
    for slide_number, coluna, titulo in pie_specs:
        idx = slide_number - 1
        if idx < 0 or idx >= len(prs.slides):
            continue
        valores = contagens_sim_nao.get(coluna, {})
        sim_value = int(valores.get("Sim", 0))
        nao_value = int(valores.get("Não", valores.get("Nao", valores.get("NÃ£o", 0))))
        slide = prs.slides[idx]
        _add_pie_block(
            slide,
            title=titulo,
            sim_value=sim_value,
            nao_value=nao_value,
            block_left=left_x,
            block_top=panel_top,
            block_w=panel_w,
            block_h=panel_h,
        )
        table_pic = None
        for shape in slide.shapes:
            if shape.shape_type == 13 and shape.left >= int(slide_w * 0.50):
                table_pic = shape

        if table_pic is not None:
            table_pic.left = right_x
            table_pic.top = panel_top
            table_pic.width = panel_w
            table_pic.height = panel_h


def gerar_ppt(
    ordered_names: list[str],
    tables_by_sheet: dict[str, SheetTable],
    arquivo_saida: str,
    assets_dir: str,
    layout_mode: str = "paired",
    contagens_sim_nao: dict[str, dict[str, int]] | None = None,
    header_layout_gap_px: int = DEFAULT_HEADER_LAYOUT_GAP_PX,
) -> None:
    assets_path = Path(assets_dir)
    primary_general = _find_primary_general_sheet(ordered_names)

    prs = Presentation()
    prs.slide_width = _inches_to_emu(13.333)
    prs.slide_height = _inches_to_emu(7.5)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    header_layout_gap_emu = int(max(0, header_layout_gap_px) * EMU_PER_PX)

    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])

    geral_slide = prs.slides.add_slide(prs.slide_layouts[6])
    top_title_h = int(slide_h * 0.20)
    _set_slide_title(geral_slide, _title_for_sheet(primary_general), slide_w, top_title_h, assets_path)

    right_half_left = int(slide_w * 0.50)
    right_half_w = int(slide_w * 0.50)
    right_margin = int(slide_w * 0.04)
    right_internal_margin = int(slide_w * 0.02)
    table_title_band_h = int(slide_h * 0.10)
    bottom_visible_gap = int(slide_h * 0.08)

    block_left = right_half_left + right_internal_margin
    block_w = right_half_w - right_margin - right_internal_margin
    block_top = top_title_h + header_layout_gap_emu
    block_h = slide_h - top_title_h - header_layout_gap_emu - bottom_visible_gap
    _add_table_block(
        geral_slide, tables_by_sheet[primary_general], block_left, block_top, block_w, block_h, slide_h, table_title_h=table_title_band_h
    )

    remaining_names = [name for name in ordered_names if name != primary_general]
    page_title_h = int(slide_h * 0.20)
    panel_top = page_title_h + header_layout_gap_emu

    single_right_half_left = int(slide_w * 0.50)
    single_right_half_w = int(slide_w * 0.50)
    single_right_margin = int(slide_w * 0.04)
    single_right_internal_margin = int(slide_w * 0.02)
    single_table_title_band_h = int(slide_h * 0.10)
    single_bottom_visible_gap = int(slide_h * 0.08)
    single_block_left = single_right_half_left + single_right_internal_margin
    single_block_w = single_right_half_w - single_right_margin - single_right_internal_margin
    single_block_top = page_title_h + header_layout_gap_emu
    single_block_h = slide_h - page_title_h - header_layout_gap_emu - single_bottom_visible_gap

    if layout_mode == "grid4":
        grid_cols = 2
        grid_rows = 2
        grid_side_margin = int(slide_w * 0.04)
        grid_col_gap = int(slide_w * 0.03)
        grid_bottom_margin = int(slide_h * 0.12)
        grid_row_gap = int(slide_h * 0.04)
        grid_total_h = slide_h - page_title_h - header_layout_gap_emu - grid_bottom_margin
        grid_panel_w = int((slide_w - (2 * grid_side_margin) - ((grid_cols - 1) * grid_col_gap)) / grid_cols)
        grid_panel_h = int((grid_total_h - ((grid_rows - 1) * grid_row_gap)) / grid_rows)

        i = 0
        while i < len(remaining_names):
            current_name = remaining_names[i]
            if _is_general_sheet(current_name):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _set_slide_title(slide, _title_for_sheet(current_name), slide_w, page_title_h, assets_path)
                _add_table_block(slide, tables_by_sheet[current_name], single_block_left, single_block_top, single_block_w, single_block_h, slide_h, table_title_h=single_table_title_band_h)
                i += 1
                continue

            chunk: list[str] = []
            while i < len(remaining_names) and len(chunk) < 4 and not _is_general_sheet(remaining_names[i]):
                chunk.append(remaining_names[i])
                i += 1
            if not chunk:
                continue

            slide_title = TITLE_P3 if any(_is_p3_sheet(name) for name in chunk) else TITLE_P1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_title(slide, slide_title, slide_w, page_title_h, assets_path)

            for idx, sheet_name in enumerate(chunk):
                row = idx // grid_cols
                col = idx % grid_cols
                x = grid_side_margin + col * (grid_panel_w + grid_col_gap)
                y = panel_top + row * (grid_panel_h + grid_row_gap)
                _add_table_block(slide, tables_by_sheet[sheet_name], x, y, grid_panel_w, grid_panel_h, slide_h)

        _add_requested_pies(
            prs,
            slide_w=slide_w,
            slide_h=slide_h,
            contagens_sim_nao=contagens_sim_nao,
            header_layout_gap_emu=header_layout_gap_emu,
        )
        prs.save(arquivo_saida)
        return

    outer_margin = int(slide_w * 0.04)
    center_gap = int(slide_w * 0.03)
    panel_w = int((slide_w - (2 * outer_margin) - center_gap) / 2)
    bottom_margin = int(slide_h * 0.15)
    panel_h = slide_h - page_title_h - header_layout_gap_emu - bottom_margin

    i = 0
    while i < len(remaining_names):
        left_name = remaining_names[i]
        if _is_general_sheet(left_name):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_title(slide, _title_for_sheet(left_name), slide_w, page_title_h, assets_path)
            _add_table_block(slide, tables_by_sheet[left_name], single_block_left, single_block_top, single_block_w, single_block_h, slide_h, table_title_h=single_table_title_band_h)
            i += 1
            continue

        right_name = remaining_names[i + 1] if (i + 1) < len(remaining_names) else None
        if right_name is not None and _is_general_sheet(right_name):
            right_name = None

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_title = _title_for_sheet(left_name)
        if right_name is not None and _is_p3_sheet(right_name):
            slide_title = TITLE_P3
        _set_slide_title(slide, slide_title, slide_w, page_title_h, assets_path)

        left_x = outer_margin
        right_x = outer_margin + panel_w + center_gap
        _add_table_block(slide, tables_by_sheet[left_name], left_x, panel_top, panel_w, panel_h, slide_h)
        if right_name is not None:
            _add_table_block(slide, tables_by_sheet[right_name], right_x, panel_top, panel_w, panel_h, slide_h)
        i += 2 if right_name is not None else 1

    _add_requested_pies(
        prs,
        slide_w=slide_w,
        slide_h=slide_h,
        contagens_sim_nao=contagens_sim_nao,
        header_layout_gap_emu=header_layout_gap_emu,
    )
    prs.save(arquivo_saida)
